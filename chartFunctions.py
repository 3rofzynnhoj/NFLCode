##functions
#####python powerpoint 
import os
import csv
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.chart.data import XyChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches
from pptx import Presentation
from pptx.chart.data import CategoryChartData,ChartData
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_CONNECTOR,MSO_SHAPE
from pptx.enum.chart import *
from pptx.util import Inches 
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN
from chartFunctions import *


from pptx.dml.color import RGBColor
import pandas as pd
import numpy as np
import logging

from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.util import Pt
 #fontipT.color.rgb = RGBColor(255, 255, 255)
#functions
 ###define colors
 
white = RGBColor(255, 255, 255)
black = RGBColor(0,0,0)
grayGap = RGBColor(116,129,129)
blue = RGBColor(2,41,200)
orange = RGBColor(254,75,56)
gray = RGBColor(173,201,211)

#slide 9 background
blue_ = RGBColor(191,220,254)
orange_ = RGBColor(255,237,230)
gray_ = RGBColor(222,233,237)

def fontChanger(shape, text, size, color, bold, align):
  text_frame = shape.text_frame
  text_frame.clear()
  p = text_frame.paragraphs[0]
  p.alignment = align
  run = p.add_run()
  run.text = text
  font = run.font
  font.bold = bold
  font.size = Pt(size)
  font.color.rgb = color
	
	
	###maka a bunch of small functions to add to this.
def fontChangerX(shape, text, size, color, bold, align, fontName):
  text_frame = shape.text_frame
  text_frame.clear()
  p = text_frame.paragraphs[0]
  p.alignment = align
  run = p.add_run()
  run.text = text
  font = run.font
  font.bold = bold
  font.size = Pt(size)
  font.color.rgb = color
  font.name = fontName
	
def fontChangerX_readIn(shape, col_name, size, color, bold, align, fontName, d, df):
  text = str(int(round(df[str(col_name)][d])))
  text_frame = shape.text_frame
  text_frame.clear()
  p = text_frame.paragraphs[0]
  p.alignment = align
  run = p.add_run()
  run.text = text
  font = run.font
  font.bold = bold
  font.size = Pt(size)
  font.color.rgb = color
  font.name = fontName

def addFont(shape, text, size, color, bold, align, fontName):
  p = shape.text_frame.paragraphs[0]
  run = p.add_run()
  run.text = text
  font = run.font
  font.bold = bold
  font.size = Pt(size)
  font.name = fontName
  font.color.rgb = color

def percentBox(shape, text):
  fontChangerX(shape, text, 120, black, True, PP_ALIGN.CENTER, "Helvetica")
  addFont(shape, "%", 44, black, True, PP_ALIGN.CENTER, "Helvetica")

def percentBox_notCenter(shape, text, align):
  fontChangerX(shape, text, 120, black, True, align, "Helvetica")
  addFont(shape, "%", 44, black, True, align, "Helvetica")	

def percentBox_notCenter_diffSize(shape, text, align, num_size, percent_size):
  fontChangerX(shape, text, num_size, black, True, align, "Helvetica")
  addFont(shape, "%", percent_size, black, True, align, "Helvetica")	

def percentBox_readIn(shape, col_name, d, df):
  text = str(int(round(df[str(col_name)][d])))
  percentBox(shape, text)
  
def percentBox_notCenter_readIn(shape, col_name, align, d, df):
  text = str(int(round(df[str(col_name)][d])))
  percentBox(shape, text)
  
def percentBox_notCenter_diffSize_readIn(shape, col_name, align, num_size, percent_size, d, df):
  text = str(int(round(df[str(col_name)][d])))
  percentBox_notCenter_diffSize(shape, text)
  	


def footerText(object, textF):
  text_frame = object.text_frame
  text_frame.clear()
  p = text_frame.paragraphs[0]
  run = p.add_run()
  run.text = '2020 Edelman Trust Barometer.'
  font = run.font
  font.bold= True
  run2=p.add_run()
  run2.text = textF 	

  
def barColors(series, col):
  for c in range(len(series[col].values)):
    fill = series[col].points[c].format.fill
    fill.solid()
    if series[col].values[c]<50:
      fill.fore_color.rgb = orange
    elif series[col].values[c] >= 50 and series[col].values[c] <=59:
      fill.fore_color.rgb = gray
    elif series[col].values[c] >= 60 and series[col].values[c] <=100:
      fill.fore_color.rgb = blue
    else:
      loggin.error('Out of range')
			
			
def bubbles(shape, num, size, fontName):
  if num <0:
    fontChangerX(shape, str(num), size, white, True, PP_ALIGN.CENTER, fontName)
    shape.fill.fore_color.rgb = black
  if num>0:
    fontChangerX(shape, "+" + str(num), size, black, True, PP_ALIGN.CENTER, fontName)
    shape.fill.fore_color.rgb = white
  if num==0:
    fontChangerX(shape, str(num), size, white, True, PP_ALIGN.CENTER, fontName)
    shape.fill.fore_color.rgb = grayGap			

def bubbles_readIn(shape, col_name, size, fontName, d, df):
  num = int(round(df[str(col_name)][d]))
  if num <0:
    fontChangerX(shape, str(num), size, white, True, PP_ALIGN.CENTER, fontName)
    shape.fill.fore_color.rgb = black
  if num>0:
    fontChangerX(shape, "+" + str(num), size, black, True, PP_ALIGN.CENTER, fontName)
    shape.fill.fore_color.rgb = white
  if num==0:
    fontChangerX(shape, str(num), size, white, True, PP_ALIGN.CENTER, fontName)
    shape.fill.fore_color.rgb = grayGap
		
def populateChart(charty, twoDArray):
  new_data = CategoryChartData()
  data_only =[]
  labels_only =[]
  for place in range(len(twoDArray)):
    data_only.append(twoDArray[place][1])
    labels_only.append(twoDArray[place][0])   
  new_data.categories = labels_only 
  new_data.add_series('Series 1', data_only)
  charty.replace_data(new_data)

def populateChart_colors(charty, twoDArray):
  new_data = CategoryChartData()
  data_only =[]
  labels_only =[]
  for place in range(len(twoDArray)):
    data_only.append(twoDArray[place][1])
    labels_only.append(twoDArray[place][0])   
  new_data.categories = labels_only 
  new_data.add_series('Series 1', data_only)
  charty.replace_data(new_data)
  series = charty.series
  barColors(series, 0)

def populateTable(tableY, twoDArray, size, color, color2, bold, bold2, align, align2, font):
  for y in range(len(twoDArray)):
    fontChangerX(tableY.cell(y,0), str(twoDArray[y][0]), size, color, bold, align, font)
    fontChangerX(tableY.cell(y,1), str(twoDArray[y][1]), size, color2, bold2, align, font)

	
def populateTable_headline(tableY, twoDArray, size, color, color2, bold, bold2, align, align2, font, headlineText, headlineSize, headlineColor, headlineBold, headlineAlign):
  fontChangerX(tableY.cell(0,0), headlineText, headlineSize, headlineColor, headlineBold, headlineAlign, font)
  for y in range(len(twoDArray)):
    d = y+1
    fontChangerX(tableY.cell(d,0), str(twoDArray[y][0]), size, color, bold, align, font)
    fontChangerX(tableY.cell(d,1), str(twoDArray[y][1]), size, color2, bold2, align2, font)

	
#####reads in a list of variable labels, and returns a 2-d array with data and labels.	
def dataImport_LabelMatch_sorted(df, variable_arr, d):
  empty_arr = []
  for g in range(len(variable_arr)):
    active_var = df[str(variable_arr[g])][d]
    empty_arr.append([variable_arr[g], active_var])
  empty_arr.sort(key = lambda x: x[0], reverse = False)
  empty_arr.sort(key = lambda x: x[1], reverse = True)
  for g in range(len(empty_arr)):
    empty_arr[g][1] = int(round(empty_arr[g][1]))
  return(empty_arr)
	
def dataImport(df, col_start, variable_arr, d):
  empty_arr = []
  for g in range(len(variable_arr)):
    col = col_start+g
    active_var = df.iloc[d][col]
    empty_arr.append([variable_arr[g], active_var])
  for g in range(len(empty_arr)):
    empty_arr[g][1] = int(round(empty_arr[g][1]))
  return(empty_arr)

	
def read_inData(df, col_name, d):
  active_var = int(round(df[str(col_name)][d]))
  return(active_var)



def flag_(slide, flag):
  pic = slide.shapes.add_picture(flag, Inches(0.0), Inches(6.97), width = Inches(.4), height = Inches(.26))		
	
	